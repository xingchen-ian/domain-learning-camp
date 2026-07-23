#!/usr/bin/env python3
"""Summer camp summary deck — matches course site art direction."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Site palette (index.html :root)
BG = RGBColor(0x09, 0x10, 0x15)
PAPER = RGBColor(0x0E, 0x15, 0x18)
INK = RGBColor(0xED, 0xF2, 0xEF)
MUTED = RGBColor(0xA8, 0xB5, 0xB1)
LEAD = RGBColor(0xD3, 0xDD, 0xDA)
EYEBROW = RGBColor(0xD9, 0xB8, 0xAD)
ACCENT = RGBColor(0xC8, 0x52, 0x43)
ACCENT_LIGHT = RGBColor(0xE2, 0x7B, 0x6D)
TEAL = RGBColor(0x78, 0xBC, 0xB6)
WARM = RGBColor(0xD2, 0xB4, 0x8B)
BLUE = RGBColor(0x86, 0xA9, 0xC9)
LINE = RGBColor(0x3A, 0x4A, 0x4E)
FOOTER = RGBColor(0x6F, 0x82, 0x83)
GRID = RGBColor(0x1A, 0x24, 0x28)

FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)
MARGIN_L = Inches(0.95)
CONTENT_W = Inches(11.4)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_grid_lines(slide):
    """Subtle vertical rhythm like the site’s repeating stripe."""
    for i in range(1, 12):
        x = Inches(0.55 + i * 1.05)
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(0), Inches(0.01), H
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRID
        line.line.fill.background()


def add_top_rule(slide):
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.03)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def add_footer(slide, page=None, total=12):
    left = slide.shapes.add_textbox(MARGIN_L, Inches(6.95), Inches(9), Inches(0.3))
    p = left.text_frame.paragraphs[0]
    p.text = "NYU Shanghai · IMA Game Design Summer Camp"
    p.font.name = FONT
    p.font.size = Pt(11)
    p.font.color.rgb = FOOTER

    if page is not None:
        right = slide.shapes.add_textbox(Inches(11.2), Inches(6.95), Inches(1.3), Inches(0.3))
        rp = right.text_frame.paragraphs[0]
        rp.text = f"{page:02d} / {total:02d}"
        rp.font.name = FONT
        rp.font.size = Pt(11)
        rp.font.color.rgb = FOOTER
        rp.alignment = PP_ALIGN.RIGHT


def textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, name=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def multiline(slide, x, y, w, h, lines, *, size=20, color=INK, gap=14, bold=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        p.line_spacing = 1.25
    return box


def section_header(slide, section, title, subtitle=None):
    textbox(
        slide, MARGIN_L, Inches(0.55), CONTENT_W, Inches(0.35),
        section.upper(), size=12, bold=True, color=EYEBROW,
    )
    textbox(
        slide, MARGIN_L, Inches(0.95), CONTENT_W, Inches(0.7),
        title, size=34 if len(title) < 40 else 30, bold=True, color=INK,
    )
    y = Inches(1.7)
    if subtitle:
        textbox(
            slide, MARGIN_L, Inches(1.65), CONTENT_W, Inches(0.4),
            subtitle, size=16, color=MUTED,
        )
        y = Inches(2.15)

    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, MARGIN_L, y, Inches(1.6), Inches(0.03)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    return y + Inches(0.35)


def card(slide, x, y, w, h, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAPER
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.04)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = accent
    top.line.fill.background()
    return shape


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    page = 0
    total = 12

    def new_slide(with_grid=True):
        nonlocal page
        page += 1
        s = prs.slides.add_slide(blank)
        set_bg(s)
        if with_grid:
            add_grid_lines(s)
        add_top_rule(s)
        return s

    # ── 1 Title ──────────────────────────────────────────────
    s = new_slide()
    textbox(
        s, MARGIN_L, Inches(2.0), CONTENT_W, Inches(0.35),
        "SUMMARY PRESENTATION", size=12, bold=True, color=EYEBROW,
    )
    title_box = s.shapes.add_textbox(
        MARGIN_L, Inches(2.4), CONTENT_W, Inches(1.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "From Real Life"
    p.font.name = FONT
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.text = "to Playable Worlds"
    p2.font.name = FONT
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = INK

    textbox(
        s, MARGIN_L, Inches(4.5), CONTENT_W, Inches(0.4),
        "NYU Shanghai · IMA Game Design Summer Camp", size=18, color=MUTED,
    )
    textbox(
        s, MARGIN_L, Inches(5.05), CONTENT_W, Inches(0.35),
        "You create. AI assists.", size=16, color=WARM,
    )
    add_footer(s, page, total)

    # ── 2 Opening quote ──────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "01 · Opening", "Opening")
    q = s.shapes.add_textbox(MARGIN_L, Inches(2.8), CONTENT_W, Inches(2.2))
    tf = q.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "A real-life experience can become\na playable game system."
    p.font.name = FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INK
    p.line_spacing = 1.2

    textbox(
        s, MARGIN_L, Inches(5.3), CONTENT_W, Inches(0.4),
        "Not a coding camp  ·  Not AI-only design", size=16, color=MUTED,
    )
    add_footer(s, page, total)

    # ── 3 Opening detail ─────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "01 · Opening", "You Are the Creator")
    multiline(
        s, MARGIN_L, content_top, CONTENT_W, Inches(3.8),
        [
            "This is not a coding camp, and we do not hand design or development over to AI.",
            "You choose the experience, the system, and the purpose.",
            "AI agents are your assistants — they build, document, and help you iterate.",
        ],
        size=22, color=LEAD, gap=22,
    )
    add_footer(s, page, total)

    # ── 4 Takeaways ──────────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "02", "What Students Take Away")
    items = [
        (TEAL, "A playable game on the web"),
        (WARM, "A way of thinking: abstract a game system from real-life experience"),
        (BLUE, "AI toolkits for recording the design and development process"),
        (ACCENT_LIGHT, "A project website that presents the design and lets others play"),
    ]
    for i, (accent, text) in enumerate(items):
        y = content_top + Inches(i * 0.95)
        card(s, MARGIN_L, y, CONTENT_W, Inches(0.8), accent)
        textbox(
            s, MARGIN_L + Inches(0.35), y + Inches(0.2),
            CONTENT_W - Inches(0.7), Inches(0.45),
            text, size=18, color=INK,
        )
    add_footer(s, page, total)

    # ── 5 Process ────────────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "03", "Design & Development Process")
    steps = [
        (TEAL, "01", "Learn from playing a game", "Not only whether it is fun — but how it teaches"),
        (WARM, "02", "Abstract an interactive system", "Clear for AI and for programming"),
        (BLUE, "03", "Build a playable demo with AI", "Implement, test, and refine"),
    ]
    card_w = Inches(3.6)
    gap = Inches(0.3)
    for i, (accent, num, title, body) in enumerate(steps):
        x = MARGIN_L + i * (card_w + gap)
        card(s, x, content_top, card_w, Inches(3.4), accent)
        textbox(
            s, x + Inches(0.3), content_top + Inches(0.35),
            card_w - Inches(0.6), Inches(0.4),
            num, size=14, bold=True, color=accent,
        )
        textbox(
            s, x + Inches(0.3), content_top + Inches(1.0),
            card_w - Inches(0.6), Inches(1.2),
            title, size=20, bold=True, color=INK,
        )
        textbox(
            s, x + Inches(0.3), content_top + Inches(2.35),
            card_w - Inches(0.6), Inches(0.8),
            body, size=15, color=MUTED,
        )
    add_footer(s, page, total)

    # ── 6 Human & AI ─────────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "04", "Human & AI: Role Distribution")
    col_w = Inches(5.5)
    # Creator
    card(s, MARGIN_L, content_top, col_w, Inches(4.0), TEAL)
    textbox(
        s, MARGIN_L + Inches(0.35), content_top + Inches(0.3),
        col_w - Inches(0.7), Inches(0.4),
        "The creator (you)", size=18, bold=True, color=INK,
    )
    multiline(
        s, MARGIN_L + Inches(0.35), content_top + Inches(0.9),
        col_w - Inches(0.7), Inches(2.9),
        [
            "Choose domain and real-life experience",
            "Turn elements into data, rules, feedback",
            "Describe and visualize the system",
            "Iterate with AI — review critically",
            "Reflect on decisions and collaboration",
        ],
        size=15, color=MUTED, gap=10,
    )
    # AI
    x2 = MARGIN_L + col_w + Inches(0.4)
    card(s, x2, content_top, col_w, Inches(4.0), WARM)
    textbox(
        s, x2 + Inches(0.35), content_top + Inches(0.3),
        col_w - Inches(0.7), Inches(0.4),
        "AI assistants", size=18, bold=True, color=INK,
    )
    multiline(
        s, x2 + Inches(0.35), content_top + Inches(0.9),
        col_w - Inches(0.7), Inches(2.9),
        [
            "Accelerate building and documentation",
            "Help implement and iterate quickly",
            "Suggest alternatives — not judgment",
            "Cannot decide why the game is worth making",
        ],
        size=15, color=MUTED, gap=10,
    )
    add_footer(s, page, total)

    # ── 7 Purpose ────────────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "04 · Human & AI", "Who Owns the Purpose?")
    q = s.shapes.add_textbox(MARGIN_L, Inches(2.9), CONTENT_W, Inches(2.4))
    tf = q.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI can speed up design and development,\nbut it cannot decide why this game is worth making."
    p.font.name = FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK
    p.line_spacing = 1.25
    add_footer(s, page, total)

    # ── 8 Value ──────────────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "05", "Why This Matters for Future Designers")
    values = [
        (WARM, "Design thinking", "Start from real life, not flashy visuals or whatever AI suggests first"),
        (TEAL, "Systems thinking", "See how data, actions, feedback, and challenges connect"),
        (BLUE, "Judgment with AI", "Know what to accept, reject, or revise"),
        (ACCENT_LIGHT, "Expression & exhibition", "Present your game and your design clearly"),
        (WARM, "Playtesting & iteration", "Test with players, learn from failure, improve"),
    ]
    for i, (accent, title, body) in enumerate(values):
        y = content_top + Inches(i * 0.78)
        bar = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, MARGIN_L, y + Inches(0.08), Inches(0.06), Inches(0.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        textbox(
            s, MARGIN_L + Inches(0.3), y,
            Inches(3.4), Inches(0.4),
            title, size=16, bold=True, color=INK,
        )
        textbox(
            s, MARGIN_L + Inches(3.8), y,
            Inches(7.5), Inches(0.55),
            body, size=15, color=MUTED,
        )
    add_footer(s, page, total)

    # ── 9 How to play ────────────────────────────────────────
    s = new_slide()
    content_top = section_header(
        s, "06", "How to Play Students’ Game Designs",
        subtitle="xingchen-ian.github.io/domain-learning-camp/gallery/",
    )
    steps = [
        "Open the Game Gallery and choose a title",
        "Read the project home page — what the game teaches",
        "Click Play and try one level",
        "Ask the student: What should the player learn? What changed after testing?",
    ]
    for i, text in enumerate(steps):
        y = content_top + Inches(i * 0.85)
        textbox(
            s, MARGIN_L, y, Inches(0.6), Inches(0.45),
            f"{i + 1:02d}", size=20, bold=True, color=TEAL,
        )
        textbox(
            s, MARGIN_L + Inches(0.7), y + Inches(0.05),
            CONTENT_W - Inches(0.7), Inches(0.5),
            text, size=18, color=LEAD,
        )
    add_footer(s, page, total)

    # ── 10 Trailers ──────────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "07", "Game Trailers")
    textbox(
        s, MARGIN_L, Inches(3.2), CONTENT_W, Inches(1.2),
        "Short clips · Screenshots · One-line pitches",
        size=24, color=MUTED,
    )
    textbox(
        s, MARGIN_L, Inches(4.4), CONTENT_W, Inches(0.5),
        "Leave this slide for live demos or embedded video.",
        size=16, color=FOOTER,
    )
    add_footer(s, page, total)

    # ── 11 Final words ───────────────────────────────────────
    s = new_slide()
    content_top = section_header(s, "08", "Final Words")
    q = s.shapes.add_textbox(MARGIN_L, Inches(2.7), CONTENT_W, Inches(3.0))
    tf = q.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "You are not here to become an “AI expert” overnight.\n\n"
        "You are here to begin thinking, crafting, and playing "
        "as a designer — with AI as a tool, not a substitute for your judgment."
    )
    p.font.name = FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = INK
    p.line_spacing = 1.3
    add_footer(s, page, total)

    # ── 12 Thank you ─────────────────────────────────────────
    s = new_slide()
    textbox(
        s, MARGIN_L, Inches(2.6), CONTENT_W, Inches(0.9),
        "Thank You", size=48, bold=True, color=INK,
    )
    textbox(
        s, MARGIN_L, Inches(3.7), CONTENT_W, Inches(0.45),
        "Play the gallery →  xingchen-ian.github.io/domain-learning-camp/gallery/",
        size=18, color=TEAL,
    )
    textbox(
        s, MARGIN_L, Inches(4.4), CONTENT_W, Inches(0.4),
        "NYU Shanghai · IMA Game Design Summer Camp", size=15, color=MUTED,
    )
    add_footer(s, page, total)

    out = Path(__file__).resolve().parent / "summer-camp-summary.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
